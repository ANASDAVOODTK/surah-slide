import streamlit as st
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io
import random
import os
from pptx.dml.color import RGBColor

# Function to call the API and fetch data
def fetch_ayahs(surah_number, range_start, range_end):
    api_url = f"https://lalithasaram.net/lalithsaram-api/ayatran/{surah_number}/{range_start}-{range_end}"
    response = requests.get(api_url)
    if response.status_code == 200:
        return response.json()
    else:
        return None

def generate_ppt(ayahs):
    prs = Presentation()
    images_path = 'images'
    image_files = [os.path.join(images_path, f) for f in os.listdir(images_path) if os.path.isfile(os.path.join(images_path, f))]

    for ayah in ayahs:
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Set slide background to a random image from the directory
        if image_files:
            bg_image = random.choice(image_files)
            slide.shapes.add_picture(bg_image, 0, 0, width=prs.slide_width, height=prs.slide_height)

        # Create a text box for the ayah text (title) at the top
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.75))
        tf_title = title_box.text_frame
        tf_title.text = ayah['ayatext']  # Set ayah text as the title
        p_title = tf_title.paragraphs[0]
        p_title.font.size = Pt(26)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(255, 255, 255)  # White font color for contrast
        p_title.alignment = PP_ALIGN.CENTER

        # Create a text box for translations and word meanings below the title
        text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        tf = text_box.text_frame
        tf.word_wrap = True

        # Set translation text, center align it, and set font color to white
        p = tf.add_paragraph()
        p.text = ayah['ayatran']
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.line_spacing = Pt(28)
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Iterate through each word for word meanings
        for word in ayah['ayawords']:
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.line_spacing = Pt(28)

            run = p.add_run()
            run.text = f"{word['meaning']} ({word['word']}) "
            run.font.color.rgb = RGBColor(255, 255, 255)

    return prs

# Streamlit UI
st.title('Generate PPT from Ayahs')

surah_number = st.number_input('Enter Surah Number', min_value=1, value=23)
range_input = st.text_input('Enter Range (e.g., 1-10)', value='1-10')
range_start, range_end = map(int, range_input.split('-'))

if st.button('Generate PPT'):
    ayahs = fetch_ayahs(surah_number, range_start, range_end)
    if ayahs:
        prs = generate_ppt(ayahs)
        ppt_file = io.BytesIO()
        prs.save(ppt_file)
        ppt_file.seek(0)
        st.download_button(label="Download PPT", data=ppt_file, file_name=f"Surah_{surah_number}_{range_start}-{range_end}.pptx", mime="application/vnd.ms-powerpoint")
    else:
        st.error('Failed to fetch data. Please check the inputs and try again.')
